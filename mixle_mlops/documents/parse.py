"""Extract text from uploaded documents and chunk it for retrieval.

Supported formats (dispatched by extension, with content-type as a fallback hint):

  * ``.txt`` / ``.md``        — decoded as UTF-8 (lenient).
  * ``.pdf``                  — ``pypdf`` (lazy import).
  * ``.docx``                 — ``python-docx`` (lazy import).
  * ``.pptx``                 — ``python-pptx`` (lazy import).

The heavy parsers are **lazy-imported inside the dispatch** so the package imports cleanly without them and a
minimal install only pays for the formats it actually uses. They are reported as the ``documents`` extra
(``pypdf``, ``python-docx``, ``python-pptx``).

Chunking is a sliding window over an approximate token estimate (``~chars/4``) with overlap, which keeps related
sentences together for embedding without a tokenizer dependency. ``chunk_text`` returns plain strings;
``parse_and_chunk`` does extract-then-chunk in one call.
"""
from __future__ import annotations

import io
from pathlib import PurePosixPath

# Rough chars-per-token used to map a token budget to a character window without a tokenizer dependency.
CHARS_PER_TOKEN = 4


class DocumentParseError(Exception):
    """Raised when a document can't be parsed (unknown/unsupported format or a missing optional parser)."""


def _ext(filename: str) -> str:
    return PurePosixPath(filename or "").suffix.lower().lstrip(".")


def _format_of(filename: str, content_type: str | None) -> str:
    """Resolve a logical format from the filename extension, falling back to the mime type."""
    ext = _ext(filename)
    if ext in {"txt", "text", "log"}:
        return "txt"
    if ext in {"md", "markdown"}:
        return "md"
    if ext == "pdf":
        return "pdf"
    if ext in {"docx"}:
        return "docx"
    if ext in {"pptx"}:
        return "pptx"
    ct = (content_type or "").split(";", 1)[0].strip().lower()
    mime_map = {
        "text/plain": "txt",
        "text/markdown": "md",
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    }
    if ct in mime_map:
        return mime_map[ct]
    if ct.startswith("text/"):
        return "txt"
    return ext or "unknown"


def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    try:
        import pypdf  # lazy: documents extra
    except ImportError as exc:  # pragma: no cover - exercised only without the extra
        raise DocumentParseError(
            "PDF support needs the 'documents' extra (pip install pypdf)."
        ) from exc
    reader = pypdf.PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    return "\n\n".join(p for p in parts if p)


def _extract_docx(data: bytes) -> str:
    try:
        import docx  # python-docx, lazy: documents extra
    except ImportError as exc:  # pragma: no cover
        raise DocumentParseError(
            "DOCX support needs the 'documents' extra (pip install python-docx)."
        ) from exc
    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text]
    for table in document.tables:                          # include table cell text
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx, lazy: documents extra
    except ImportError as exc:  # pragma: no cover
        raise DocumentParseError(
            "PPTX support needs the 'documents' extra (pip install python-pptx)."
        ) from exc
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def extract_text(data: bytes, *, filename: str = "", content_type: str | None = None) -> str:
    """Extract plain text from ``data`` for a supported format (dispatched by extension/mime)."""
    fmt = _format_of(filename, content_type)
    if fmt in {"txt", "md"}:
        return _extract_txt(data)
    if fmt == "pdf":
        return _extract_pdf(data)
    if fmt == "docx":
        return _extract_docx(data)
    if fmt == "pptx":
        return _extract_pptx(data)
    raise DocumentParseError(
        f"unsupported document format {fmt!r} (filename={filename!r}, content_type={content_type!r})"
    )


def chunk_text(
    text: str,
    *,
    chunk_tokens: int = 256,
    overlap_tokens: int = 32,
    chunk_chars: int | None = None,
    overlap_chars: int | None = None,
) -> list[str]:
    """Sliding-window chunking with overlap.

    Sizes are given in approximate tokens (mapped to characters via :data:`CHARS_PER_TOKEN`); pass
    ``chunk_chars`` / ``overlap_chars`` to specify a character window directly. Windows are cut at the nearest
    whitespace before the boundary when possible, so chunks don't split mid-word.
    """
    text = (text or "").strip()
    if not text:
        return []
    size = chunk_chars if chunk_chars is not None else max(1, chunk_tokens * CHARS_PER_TOKEN)
    over = overlap_chars if overlap_chars is not None else max(0, overlap_tokens * CHARS_PER_TOKEN)
    over = min(over, size - 1)                              # overlap must be strictly less than the window
    chunks: list[str] = []
    n = len(text)
    start = 0
    while start < n:
        end = min(start + size, n)
        if end < n:                                        # try to break on whitespace for cleaner chunks
            window = text[start:end]
            cut = window.rfind(" ")
            if cut > size // 2:                            # only honour the break if it isn't too early
                end = start + cut
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(end - over, start + 1)
    return chunks


def parse_and_chunk(
    data: bytes,
    *,
    filename: str = "",
    content_type: str | None = None,
    chunk_tokens: int = 256,
    overlap_tokens: int = 32,
) -> tuple[str, list[str]]:
    """Extract text then chunk it. Returns ``(full_text, chunks)``."""
    text = extract_text(data, filename=filename, content_type=content_type)
    chunks = chunk_text(text, chunk_tokens=chunk_tokens, overlap_tokens=overlap_tokens)
    return text, chunks
